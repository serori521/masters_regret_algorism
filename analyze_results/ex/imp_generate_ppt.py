"""
PowerPoint専用ヒートマップ発表資料生成ツール（5枚/スライド・ルール別・軽量版）

前提:
    pip install python-pptx pillow

想定ファイル構造:
    analyze_the_results/heatmaps/<rule>/u{1,2}/N6/heatmap_<rule>_<method>_u{u}_N6_{TW}.png
"""

import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image  # ★ 画像圧縮用


def collect_images_for_rule(rule: str,
                            image_root: str = "analyze_the_results/heatmaps") -> list:
    pattern = os.path.join(image_root, rule, "u*", "N6", "*.png")
    png_files = glob.glob(pattern)
    png_files.sort()
    return png_files


def parse_filename_info(filename: str):
    base = os.path.basename(filename)
    parts = base.replace(".png", "").split("_")

    method = "UNKNOWN"
    utility = "u?"
    n_val = "N?"
    true_weight = "?"

    if len(parts) >= 6:
        # parts[0] = "heatmap"
        # parts[1] = rule
        method = parts[2]
        utility = parts[3]
        n_val = parts[4]
        true_weight = parts[5]

    return method, utility, n_val, true_weight


# ========= 画像圧縮まわり =========

def compress_image(src_path: str,
                   dst_dir: str,
                   max_size: int = 1200,
                   quality: int = 60) -> str:
    """
    PNG などの元画像を「軽量 JPEG」に変換して保存し、そのパスを返す。

    max_size: 縦横の最大ピクセル数（長辺がこれを超えたら縮小）
    quality: JPEG 品質（60〜70くらいで 1/5〜1/10 になることも多い）
    """
    os.makedirs(dst_dir, exist_ok=True)

    base = os.path.basename(src_path)
    name_no_ext = os.path.splitext(base)[0]
    dst_path = os.path.join(dst_dir, f"{name_no_ext}_q{quality}.jpg")

    # すでに作っているなら再利用
    if os.path.exists(dst_path):
        return dst_path

    try:
        img = Image.open(src_path)
        img = img.convert("RGB")  # PNG(透過)でもJPEGにするのでRGBへ

        w, h = img.size
        scale = min(1.0, max_size / max(w, h))
        if scale < 1.0:
            new_w = int(w * scale)
            new_h = int(h * scale)
            img = img.resize((new_w, new_h), Image.LANCZOS)

        img.save(dst_path, format="JPEG", quality=quality, optimize=True)
        return dst_path
    except Exception as e:
        print(f"[WARN] 画像圧縮に失敗: {src_path} -> {e}")
        # 失敗したら元のパスを返す（最悪そのまま埋め込む）
        return src_path


def create_five_per_slide_for_rule(
    rule: str,
    image_root: str = "analyze_the_results/heatmaps",
    output_dir: str = "/workspaces/inulab_julia_devcontainer/results/pptx",
    compressed_dir: str = "/workspaces/inulab_julia_devcontainer/results/pptx/_compressed_images",
    max_size: int = 1200,
    quality: int = 60,
) -> bool:

    print(f"\n📊 5枚/スライド PowerPoint 生成（method別・軽量版）: rule = {rule}")

    png_files = collect_images_for_rule(rule, image_root=image_root)
    if not png_files:
        print(f"❌ 画像ファイルが見つかりません: {image_root}/{rule}/u*/N6/*.png")
        return False

    os.makedirs(output_dir, exist_ok=True)
    output_file = os.path.join(output_dir, f"{rule}_heatmaps_5per_method_light.pptx")

    prs = Presentation()

    # タイトル
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"{rule.capitalize()} Heatmaps (Grouped by Method, Compressed)"
    slide.placeholders[1].text = (
        f"Total images: {len(png_files)}\n"
        f"max_size={max_size}, JPEG quality={quality}"
    )

    # (method, utility, N6) ごとにグループ化
    groups = {}   # key = (method, utility, n_val), value = [file list]

    for file in png_files:
        method, utility, n_val, tw = parse_filename_info(file)
        key = (method, utility, n_val)
        groups.setdefault(key, []).append((tw, file))

    # TW順にソートしてスライド生成
    for key, items in groups.items():
        method, utility, n_val = key

        items_sorted = sorted(items, key=lambda x: x[0])  # x = (TW, filename)
        png_list = [f for (_, f) in items_sorted]

        for i in range(0, len(png_list), 5):
            batch = png_list[i:i+5]

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

            title_text = f"{rule.upper()} - {method} ({utility}, {n_val})"
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(11), Inches(0.6))
            text_frame = textbox.text_frame
            text_frame.text = title_text
            p0 = text_frame.paragraphs[0]
            p0.font.size = Pt(20)
            p0.font.bold = True
            p0.alignment = PP_ALIGN.CENTER

            positions = [
                (Inches(0.5),  Inches(1.0)),
                (Inches(4.0),  Inches(1.0)),
                (Inches(7.5),  Inches(1.0)),
                (Inches(2.25), Inches(4.2)),
                (Inches(5.75), Inches(4.2)),
            ]
            img_width = Inches(3.2)

            for j, png_file in enumerate(batch):
                if j >= len(positions):
                    break

                # ★ ここで軽量JPEGに変換してから貼り付ける
                compressed_path = compress_image(
                    png_file,
                    dst_dir=os.path.join(compressed_dir, rule),
                    max_size=max_size,
                    quality=quality,
                )

                left, top = positions[j]
                slide.shapes.add_picture(compressed_path, left, top, width=img_width)

                # ラベル
                fname = os.path.basename(png_file).replace(".png", "")
                parts = fname.split("_")
                tw = parts[-1] if parts else "?"
                label_text = f"{method}-{utility}-{n_val}-{tw}"

                label_top = top + Inches(2.5)
                box = slide.shapes.add_textbox(left, label_top, img_width, Inches(0.3))
                p = box.text_frame.paragraphs[0]
                p.text = label_text
                p.font.size = Pt(10)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

    prs.save(output_file)
    print(f"✅ 保存完了: {output_file}")
    return True


def main():
    """
    regret / maximin / Maximax の3種類について、
    それぞれ 1スライド5枚版の「軽量 PPTX」を生成する。
    """
    print("🎯 PowerPoint 5枚/スライド版ヒートマップ生成 (rule別・軽量版)")
    print("=" * 60)

    rules = ["regret", "maximin", "Maximax"]

    for rule in rules:
        # ここで max_size や quality を調整してサイズと画質のバランスを取る
        create_five_per_slide_for_rule(
            rule,
            max_size=1200,   # 例: 長辺 1200px まで
            quality=60,      # 例: JPEG 品質 60
        )

    print("=" * 60)
    print("✅ すべての rule の 軽量PPTX 生成処理が終了しました。")


if __name__ == "__main__":
    main()
